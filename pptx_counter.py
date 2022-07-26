from pptx import Presentation
from os import walk

def pptx_files(path: str):
    filenames = next(walk(path), (None, None, []))[2]    
    pptx_files = [file for file in filenames if file[-4:] == 'pptx']
    return pptx_files

def counter(file: str):
    prs = Presentation(file)
    total_symbols = 0
    total_symbols_with_space = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    res = sum(not chr.isspace() for chr in run.text)     # без пробелов
                    total_symbols += res
                    total_symbols_with_space += len(run.text)
                    # print(res, run.text)
    return total_symbols, total_symbols_with_space


total = 0
total_with_space = 0
print('----------------------------------------------------------------------------------------------------------------------------------------------')
for i in pptx_files('.'):
    without_spaces, with_spaces = counter(i)
    total += without_spaces
    total_with_space += with_spaces
    print('{:<90} {:<11} {:<10}{}{:<10}{}'.format(i, ' | symbols:', with_spaces, ' | without spaces:', without_spaces, '|'))

    print('----------------------------------------------------------------------------------------------------------------------------------------------')


print('Total symbols: ', total_with_space)
print('Total without spaces: ', total)
input()



