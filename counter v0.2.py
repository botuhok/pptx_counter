import xlrd as xl                             #Import xlrd package
from openpyxl import load_workbook
from pptx import Presentation
from os import walk
from colorama import init
from termcolor import colored

path = '.'

def xls_files(path: str):
    filenames = next(walk(path), (None, None, []))[2]
    xls_files = [file for file in filenames if file[-3:] == 'xls']
    sum_all = []
    for file in xls_files:
        counter = 0
        counter_noSpace = 0
        wb = xl.open_workbook(file)
        for sheet in wb.sheets():
            for rownum in range(sheet.nrows):
                row = sheet.row_values(rownum)
                for c_el in row:
                    if c_el:
                        cell = str(c_el)
                        if cell[-2:] == '.0':
                            cell = cell[:-2]
                        counter += len(cell)
                        counter_noSpace += sum(not chr.isspace() for chr in cell)
        sum_all += (file, counter, counter_noSpace)
    return sum_all


def pptx_files(path: str):
    filenames = next(walk(path), (None, None, []))[2]
    pptx_files = [file for file in filenames if file[-4:] == 'pptx']
    sum_all = {}
    for file in pptx_files:
        prs = Presentation(file)
        counter = 0
        counter_noSpace = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        res = sum(not chr.isspace() for chr in run.text)     # без пробелов
                        counter_noSpace += res
                        counter += len(run.text)
        sum_all[file] = (counter, counter_noSpace)
    return sum_all

def xlsx_files(path: str):
    filenames = next(walk(path), (None, None, []))[2]                    # список всех файлов рядом со скриптом
    xlsx_files = [file for file in filenames if file[-4:] == 'xlsx']     # список только xlsx
    sum_all = {}
    for file in xlsx_files:
        counter = 0
        counter_noSpace = 0
        wb = load_workbook(file)
        for sheet in wb.worksheets:
            for line in sheet.values:
                for cell in line:
                    if cell:
                        cell = str(cell)
                        counter += len(cell)
                        counter_noSpace += sum(not chr.isspace() for chr in cell)
        sum_all[file] = (counter, counter_noSpace)
    return sum_all

def xls_files(path: str):
    filenames = next(walk(path), (None, None, []))[2]
    xls_files = [file for file in filenames if file[-3:] == 'xls']
    sum_all = {}
    for file in xls_files:
        counter = 0
        counter_noSpace = 0
        wb = xl.open_workbook(file)
        for sheet in wb.sheets():
            for rownum in range(sheet.nrows):
                row = sheet.row_values(rownum)
                for c_el in row:
                    if c_el:
                        cell = str(c_el)
                        if cell[-2:] == '.0':
                            cell = cell[:-2]
                        counter += len(cell)
                        counter_noSpace += sum(not chr.isspace() for chr in cell)
        sum_all[file] = (counter, counter_noSpace)
    return sum_all

sum_all = xls_files(path)
sum_all.update(xlsx_files(path))
sum_all.update(pptx_files(path))

total = 0
total_with_space = 0
init()
print()
print(colored('===============================================================================================================================', 'cyan'))
print('{:<92}{}{:<15}{}{:<10}{}'.format('ИМЯ ФАЙЛА', '|  ', 'СИМВОЛЫ', "| ","БЕЗ ПРОБЕЛОВ","  |"))
print(colored('===============================================================================================================================', 'cyan'))
for file, value in sum_all.items():
    total += value[1]
    total_with_space += value[0]
    print('{:<90} {} {:<13} {} {:<13}{}'.format(file, ' | ', value[0], ' | ', value[1], '|'))

    print(colored('-------------------------------------------------------------------------------------------------------------------------------', 'cyan'))

print()
print(colored('ИТОГО:              ', 'green'), colored(total_with_space, 'yellow'))
print()
print(colored('ИТОГО БЕЗ ПРОБЕЛОВ: ', 'green'), colored(total, 'yellow'))
input()



